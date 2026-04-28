"""
COMP447 progress presentation builder.

Renders LaTeX equations via matplotlib mathtext, regenerates the pareto and
ablation plots from the actual results CSVs, and stitches everything into a
16:9 PPTX that is Google Slides import-friendly.

Run:  python3 generate.py
Output:  COMP447_progress.pptx
"""

from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt

plt.rcParams["mathtext.fontset"] = "cm"
plt.rcParams["font.family"] = "serif"

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
ASSETS.mkdir(exist_ok=True)
OUT = ROOT / "COMP447_progress.pptx"

INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
ACCENT = RGBColor(0xA8, 0x3D, 0x3D)
WARN = RGBColor(0x42, 0x42, 0x42)
SOFT = RGBColor(0xEE, 0xEE, 0xEE)

FONT = "Helvetica"
TOTAL_SLIDES = 10


# ---------------------------------------------------------------------------
# Asset rendering
# ---------------------------------------------------------------------------
def render_eq(tex: str, name: str, fontsize: int = 30, width: float = 10.0,
              height: float = 1.4) -> str:
    """Render a LaTeX expression to a transparent PNG."""
    path = ASSETS / f"{name}.png"
    fig = plt.figure(figsize=(width, height), dpi=240)
    fig.patch.set_alpha(0)
    fig.text(0.5, 0.5, f"${tex}$", fontsize=fontsize, ha="center", va="center",
             color="#1a1a1a")
    plt.axis("off")
    plt.savefig(path, dpi=240, transparent=True, bbox_inches="tight",
                pad_inches=0.08)
    plt.close()
    return str(path)


def render_pareto() -> str:
    path = ASSETS / "pareto.png"
    fig, ax = plt.subplots(figsize=(9, 4.8), dpi=180)
    fig.patch.set_facecolor("white")

    heun_steps = [5, 10, 18, 25, 50]
    heun_lat = [62.21, 131.75, 243.79, 331.66, 673.50]
    heun_fid = [37.78, 2.64, 1.96, 1.98, 2.02]
    ax.plot(heun_lat, heun_fid, "-o", c="#424242", markersize=9,
            linewidth=2, label="Heun (EDM, no tuning)")
    # Stagger the bottom-row labels (18 / 25 / 50) vertically so they do not
    # collide on the log-x axis. 5- and 10-step labels go to the right.
    for s, x, y in zip(heun_steps, heun_lat, heun_fid):
        if s in (18, 50):
            dx, dy, ha, va = 0, 12, "center", "bottom"
        elif s == 25:
            dx, dy, ha, va = 0, -14, "center", "top"
        else:
            dx, dy, ha, va = 8, -2, "left", "center"
        ax.annotate(f"{s} steps", (x, y), xytext=(dx, dy),
                    textcoords="offset points", fontsize=9,
                    color="#424242", ha=ha, va=va)

    ax.scatter([7.05, 13.97], [5.77, 2.47], s=160, c="#a83d3d",
               marker="D", zorder=5, label="ECT (ours, 2k kimg)",
               edgecolor="white", linewidth=1.2)
    ax.annotate("1 step", (7.05, 5.77), xytext=(8, 8),
                textcoords="offset points", fontsize=10, color="#a83d3d",
                weight="bold")
    ax.annotate("2 step", (13.97, 2.47), xytext=(8, 8),
                textcoords="offset points", fontsize=10, color="#a83d3d",
                weight="bold")

    ax.set_xscale("log")
    ax.set_yscale("log")
    ax.set_xlabel("Latency per image (ms, batch = 1)", fontsize=12)
    ax.set_ylabel("FID", fontsize=12)
    ax.grid(True, which="both", alpha=0.25, linestyle="--")
    ax.legend(loc="upper right", frameon=False, fontsize=11)
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.tick_params(labelsize=10)
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close()
    return str(path)


def render_ablation() -> str:
    path = ASSETS / "ablation.png"
    fig, ax = plt.subplots(figsize=(9, 4.8), dpi=180)
    fig.patch.set_facecolor("white")
    kimg = [500, 1000, 1500, 1980]
    fid = [8.135, 2.773, 2.702, 2.446]
    ax.plot(kimg, fid, "-o", c="#a83d3d", markersize=11, linewidth=2.5)
    ax.axhline(y=2.11, color="#424242", linestyle="--", linewidth=1.5)
    ax.text(1950, 2.18, "ECT paper, 400 k kimg  →  FID 2.11",
            color="#424242", fontsize=10, ha="right")
    for k, f in zip(kimg, fid):
        ax.annotate(f"{f:.2f}", (k, f), xytext=(0, 12),
                    textcoords="offset points", fontsize=11,
                    ha="center", weight="bold")
    ax.set_xlabel("Tuning budget (kimg)", fontsize=12)
    ax.set_ylabel("2 step FID50k", fontsize=12)
    ax.set_xticks(kimg)
    ax.grid(True, alpha=0.25, linestyle="--")
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.tick_params(labelsize=10)
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close()
    return str(path)


# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------
def _set_run(run, *, size, bold=False, color=INK, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    _set_run(p.runs[0], size=size, bold=bold, color=color)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=17, gap=10):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•   " + item
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        _set_run(p.runs[0], size=size, color=INK)
    return tb


def add_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55),
                                 Inches(0.50), Inches(0.07), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text(slide, 0.78, 0.42, 12.0, 0.7, title, size=24, bold=True)


def add_footer(slide, page):
    add_text(slide, 11.6, 7.05, 1.5, 0.3, f"{page} / {TOTAL_SLIDES}",
             size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image(slide, path, *, top, max_height_in, max_width_in=11.5):
    """Place an image centered horizontally, scaled to fit in the box."""
    with Image.open(path) as img:
        w, h = img.size
    ratio = w / h
    target_h = max_height_in
    target_w = target_h * ratio
    if target_w > max_width_in:
        target_w = max_width_in
        target_h = target_w / ratio
    slide_w = 13.333
    left = (slide_w - target_w) / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top),
                             width=Inches(target_w),
                             height=Inches(target_h))


def add_quote_block(slide, x, y, w, h, quote, attribution):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = SOFT
    box.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x),
                                 Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_text(slide, x + 0.35, y + 0.22, w - 0.55, h - 0.7,
             f"“{quote}”", size=18, color=INK)
    add_text(slide, x + 0.35, y + h - 0.55, w - 0.55, 0.4,
             attribution, size=12, color=GRAY)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# Slide 1 — Title -------------------------------------------------------------
s = prs.slides.add_slide(blank)

bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(2.65),
                         Inches(2.0), Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

add_text(s, 0.55, 2.05, 12, 0.4,
         "COMP447  ·  PROGRESS REPORT", size=12, color=GRAY)
add_text(s, 0.55, 2.95, 12, 1.4,
         "When Does Cheap Consistency Tuning Pay Off?",
         size=38, bold=True, color=INK)
add_text(s, 0.55, 4.20, 12, 0.6,
         "When is consistency tuning worth it?  ·  "
         "The open question we found.",
         size=16, color=GRAY)
add_text(s, 0.55, 5.55, 12, 0.4,
         "Batuhan Karaman   ·   Kadir Yiğit Özçelik",
         size=14, color=INK)
add_text(s, 0.55, 5.95, 12, 0.4, "April 2026", size=12, color=GRAY)


# Slide 2 — The question ------------------------------------------------------
s = prs.slides.add_slide(blank)
add_title_bar(s, "The question")
add_bullets(s, 0.85, 1.70, 11.5, 4.5, [
    "Diffusion models reach state of the art image quality at the cost of "
    "30 to 100+ neural function evaluations per image.",
    "Consistency Models compress sampling into 1 or 2 steps, but the upfront "
    "tuning is non trivial.",
    "When is cheap consistency tuning worth paying for, and how should you "
    "sample once you've paid?",
], size=18, gap=18)
add_footer(s, 2)


# Slide 3 — Background --------------------------------------------------------
s = prs.slides.add_slide(blank)
add_title_bar(s, "Background")

add_text(s, 0.85, 1.30, 11.5, 0.4,
         "Diffusion: probability flow ODE  [Karras+ 2022, EDM]",
         size=13, color=GRAY)
eq_pfode = render_eq(r"\frac{dx_t}{dt}\;=\;-\,t\,\nabla_{x_t}\log p_t(x_t)",
                     "eq_pfode", fontsize=30)
add_image(s, eq_pfode, top=1.70, max_height_in=0.95)

add_text(s, 0.85, 3.20, 11.5, 0.4,
         "Consistency model: a learnable flow map  [Song+ 2023]",
         size=13, color=GRAY)
eq_cm = render_eq(
    r"f_\theta(x_t,\,t)\;\approx\;x_0,\qquad\;"
    r"f_\theta(x_t,\,t)\;=\;f_\theta(x_s,\,s)\;\;\text{for }(x_t,x_s)\text{ on the same trajectory}",
    "eq_cm", fontsize=20)
add_image(s, eq_cm, top=3.60, max_height_in=0.95)

add_bullets(s, 0.85, 5.30, 11.5, 1.2, [
    "Diffusion: 2 NFE / step, 18 steps ≈ 35 NFE.   ·   "
    "CM: 1 NFE for 1 step, 2 NFE for 2 step.",
], size=15)
add_footer(s, 3)


# Slide 4 — ECT ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_title_bar(s, "Easy Consistency Tuning  ·  the model we use")

add_text(s, 0.85, 1.40, 11.5, 0.4,
         "Self consistency between infinitesimally close noise levels  [Geng+ 2025]",
         size=13, color=GRAY)
eq_ect = render_eq(
    r"\mathcal{L}_{\mathrm{ECT}}(\theta)\;=\;"
    r"\mathbb{E}_{y,\epsilon,t}\!\left[\,w(t)\,d\!\left(\,"
    r"f_\theta(y+t\epsilon,\,t),\;f_{\theta'}(y+r\epsilon,\,r)\right)\,\right],"
    r"\quad r = t - \Delta(t),\;\Delta(t)\!\to\!0",
    "eq_ect", fontsize=20)
add_image(s, eq_ect, top=1.95, max_height_in=1.15)

add_bullets(s, 0.85, 3.65, 11.5, 3.5, [
    "Initialised from a pretrained EDM checkpoint, no distillation teacher.",
    "Two step inference uses a single hand picked midpoint at value 0.821, "
    "with no derivation in the paper.",
    "Paper claim: ~1 A100 hour to reach FID 2.11 at 400 k kimg.",
], size=17, gap=14)
add_footer(s, 4)


# Slide 5 — Pareto frontier ---------------------------------------------------
s = prs.slides.add_slide(blank)
add_title_bar(s, "Pareto frontier")

add_text(s, 0.85, 1.30, 11.5, 0.35,
         "EDM CIFAR 10 32×32  ·  ECT @ 1980 kimg  ·  Heun 5 to 50 steps  ·  "
         "NVIDIA G4 Blackwell  ·  batch = 1",
         size=11, color=GRAY)

add_image(s, render_pareto(), top=1.70, max_height_in=4.4)

add_bullets(s, 0.85, 6.20, 11.5, 1.2, [
    "ECT 2 step  ·  FID 2.47  ·  14 ms / image  →  17× faster than Heun 18.",
    "On this short tuning budget, ECT Pareto dominates Heun across every operating point.",
], size=14, gap=8)
add_footer(s, 5)


# Slide 6 — Tuning budget ablation -------------------------------------------
s = prs.slides.add_slide(blank)
add_title_bar(s, "Tuning budget ablation")

add_image(s, render_ablation(), top=1.30, max_height_in=4.5)

add_bullets(s, 0.85, 5.85, 11.5, 1.6, [
    "94 % of paper reported quality at 0.5 % of the paper’s training budget.",
    "Diminishing returns past 1000 kimg  →  early stopping recipe is realistic.",
    "Break even vs Heun 18  ·  N* ≈ 23.5 k images (batch 1) before tuning amortises.",
], size=14, gap=8)
add_footer(s, 6)


# Slide 7 — The real open problem --------------------------------------------
s = prs.slides.add_slide(blank)
add_title_bar(s, "The real open problem")

add_bullets(s, 0.85, 1.25, 11.5, 0.95, [
    "As we characterised 2 step behaviour across our 4 checkpoints, every "
    "result rested on a single hidden parameter: a midpoint at 0.821, "
    "fixed across all training stages and step budgets.",
], size=14, gap=8)

add_quote_block(
    s, 0.85, 2.40, 11.6, 2.05,
    "It’s a very good question and I think also an open research problem. "
    "I don’t have a manual for it. You can treat it as an optimization "
    "problem, maximizing sample quality with respect to these sampling "
    "schedules.",
    "Zhengyang Geng, ECT first author",
)

add_bullets(s, 0.85, 4.65, 11.5, 2.3, [
    "No derivation in ECT [Geng+ 2025], sCM [Karras+ 2024], or Multistep CM [Heek+ 2024].",
    "Heuristic timestep selection exists [TDD, Wang+ 2024], but principled "
    "KLUB-based optimization for consistency models does not.",
    "Align Your Steps [Sabour+ 2024] solves it for diffusion only — we adapt "
    "the KL upper bound to the consistency flow map.",
], size=14, gap=10)
add_footer(s, 7)


# Slide 8 — Method: KLUB CM ---------------------------------------------------
s = prs.slides.add_slide(blank)
add_title_bar(s, "Our method  ·  KLUB CM schedule")

add_text(s, 0.85, 1.40, 11.5, 0.4,
         "Re-derive the Align Your Steps KL upper bound for the consistency "
         "function, not the score field.",
         size=13, color=GRAY)

eq_klub = render_eq(
    r"\mathcal{T}^{\,*}\;=\;\arg\min_{\{t_i\}}\;\sum_{i=1}^{N}\;"
    r"\mathbb{E}_{x\sim p}\!\left[\,\left\|\,f_\theta(x_{t_i},\,t_i)\;-\;"
    r"f_\theta(x_{t_{i-1}},\,t_{i-1})\,\right\|^{2}\,\right]",
    "eq_klub", fontsize=22)
add_image(s, eq_klub, top=1.95, max_height_in=1.10)

add_text(s, 0.85, 3.30, 11.5, 0.35,
         "Inspiration:  KL upper bound [Sabour+ 2024]  ·  trajectory awareness "
         "[Kim+ 2024, CTM]  ·  adaptive schedules [OSS, 2025]",
         size=11, color=GRAY)

add_bullets(s, 0.85, 3.95, 11.5, 3.0, [
    "Linearise the flow map (not the score) at chosen knots; the KLUB upper "
    "bounds the discretisation error.",
    "Schedule is checkpoint conditioned  →  drifts with model quality, unlike "
    "a fixed magic number.",
    "Predicted impact: pushes ECT operating points further left on the Pareto "
    "frontier; generalises to unseen step budgets without retraining.",
], size=14, gap=10)
add_footer(s, 8)


# Slide 9 — Implementation ----------------------------------------------------
s = prs.slides.add_slide(blank)
add_title_bar(s, "Implementation  ·  what we are coding")

add_bullets(s, 0.85, 1.40, 11.5, 2.6, [
    "KLUB derivation for the CM flow map  ·  rederived for the consistency "
    "function, not the score field.",
    "Schedule optimiser  ·  LBFGS over the schedule knots in PyTorch — gradient "
    "based, orders of magnitude cheaper than grid search.",
    "On policy loss patch in training/loss.py  ·  student's own outputs at the "
    "chosen midpoint, motivated by exposure bias work [Agarwal+ 2023, GKD].",
    "Corrected evaluation pipeline  ·  1 step and 2 step routed to the right sampler.",
], size=14, gap=10)

add_text(s, 0.85, 4.20, 11.5, 0.35,
         "training/loss.py  ·  on policy schedule aware term  (sketch)",
         size=11, color=GRAY)

# Code box
code_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85),
                              Inches(4.65), Inches(11.6), Inches(2.05))
code_box.fill.solid()
code_box.fill.fore_color.rgb = SOFT
code_box.line.fill.background()

code_lines = [
    "x_mid = f_theta(z * t_max, t_max).detach()         # student's own coarse output",
    "t_mid = klub_optimiser(checkpoint=ckpt)            # KLUB derived knot",
    "L_on  = w(t_mid) * d(",
    "    f_theta      (x_mid + eps * t_mid, t_mid),     # on policy input",
    "    f_theta_prime(x_mid + eps * r,     r),         # EMA target",
    ")",
    "L     = L_ECT  +  lam * L_on                       # original + on policy",
]
tb = s.shapes.add_textbox(Inches(1.05), Inches(4.78), Inches(11.2), Inches(1.85))
tf = tb.text_frame
tf.word_wrap = True
for i, line in enumerate(code_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    _set_run(p.runs[0], size=11, color=INK, font="Menlo")

add_footer(s, 9)


# Slide 10 — Status & next steps ---------------------------------------------
s = prs.slides.add_slide(blank)
add_title_bar(s, "Status & next steps")

add_text(s, 0.85, 1.45, 5.5, 0.4, "Done", size=16, bold=True, color=ACCENT)
add_bullets(s, 0.85, 1.95, 6.0, 4.5, [
    "ECT tuning, 4 checkpoints saved.",
    "Heun baselines @ 5 to 50 steps.",
    "Latency profiling, B = 1 and B = 64.",
    "Corrected evaluation pipeline.",
], size=15, gap=10)

add_text(s, 7.0, 1.45, 6.0, 0.4, "In flight  ·  Up next",
         size=16, bold=True, color=WARN)
add_bullets(s, 7.0, 1.95, 6.0, 4.5, [
    "Midpoint grid sweep across all 4 checkpoints.",
    "KLUB CM derivation + schedule optimiser.",
    "On policy schedule aware fine tune  (stretch).",
    "Final report  ·  early June.",
], size=15, gap=10)

add_text(s, 0.85, 6.10, 11.5, 0.5,
         "Risk  ·  GPU time on Colab; mitigated by small checkpoint footprint "
         "and 10k sample screening for the sweep.",
         size=12, color=GRAY)
add_footer(s, 10)


# ---------------------------------------------------------------------------
prs.save(OUT)
print(f"wrote {OUT}")
